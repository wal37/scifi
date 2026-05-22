from __future__ import annotations

import sys
from pathlib import Path


BASE_DIR = Path(__file__).resolve().parent
for candidate in [
    Path("/Users/wale/Desktop/softtttt/aiapis/apiweb/.venv/lib/python3.13/site-packages"),
    Path("/Users/wale/Desktop/softtttt/foodup/foodweb/.venv/lib/python3.13/site-packages"),
]:
    if candidate.exists() and str(candidate) not in sys.path:
        sys.path.append(str(candidate))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = BASE_DIR / "LOOPNODE_PITCH_DECK.pptx"

LOGO = BASE_DIR / "loopnode-logo.png"
IMG_GAMES = BASE_DIR / "lgames.jpg"
IMG_AR = BASE_DIR / "lar.jpg"
IMG_VR = BASE_DIR / "lvr.jpg"
IMG_CONFIG = BASE_DIR / "lconfig.jpg"
IMG_ADS = BASE_DIR / "lad.jpg"


def rgb(code: str) -> RGBColor:
    code = code.replace("#", "")
    return RGBColor(int(code[0:2], 16), int(code[2:4], 16), int(code[4:6], 16))


INK = rgb("#07111E")
NAVY = rgb("#0D1726")
SLATE = rgb("#10233B")
BLUE = rgb("#2C6BFF")
CYAN = rgb("#6EDBFF")
LIME = rgb("#B8FF68")
VIOLET = rgb("#A58EFF")
AMBER = rgb("#FFD166")
PAPER = rgb("#F4F7FB")
WHITE = rgb("#FFFFFF")
MUTED = rgb("#6A7C90")
LINE = rgb("#D8E1EC")
DARK_LINE = rgb("#1E3552")


SLIDES = [
    {
        "kind": "cover",
        "eyebrow": "Browser-native 3D infrastructure",
        "title": "Loopnode builds a premium operating layer for realtime 3D on the web.",
        "subtitle": "A browser-native platform for scene creation, runtime delivery, collaborative review, and high-fidelity launch surfaces across games, XR, commerce, and interactive product experiences.",
        "image": IMG_CONFIG,
    },
    {
        "kind": "statement",
        "eyebrow": "01  Product thesis",
        "title": "Most teams still piece together modern 3D delivery from disconnected tools.",
        "body": "Scene editors, runtime code, review links, asset pipelines, deployment workflows, and analytics are usually managed as separate systems. Loopnode compresses that sprawl into one cleaner product surface.",
        "points": [
            "One browser workspace instead of fragmented tooling",
            "Faster iteration from scene edits to published output",
            "A better fit for product teams than engine-first workflows",
        ],
    },
    {
        "kind": "three_cards",
        "eyebrow": "02  Market shift",
        "title": "Why the category matters now",
        "subtitle": "Interactive 3D is moving beyond game studios into product, commerce, simulation, and immersive web experiences.",
        "cards": [
            ("Browser expectations are rising", "Premium launches now need richer visual surfaces, motion, and spatial interaction directly on the web."),
            ("WebGPU expands the ceiling", "Modern browser graphics make higher-fidelity realtime rendering practical for a wider product set."),
            ("Teams need productized infrastructure", "Creative, engineering, and go-to-market teams need a system that feels deployable, not improvised."),
        ],
    },
    {
        "kind": "image_split",
        "eyebrow": "03  Workflow",
        "title": "Loopnode keeps composition, review, and launch inside one operating flow.",
        "body": "The platform is built for teams that want design and technical control without losing speed. Scene building, logic, previewing, and deployment stay connected.",
        "image": IMG_GAMES,
        "bullets": [
            "Browser scene editing with premium output paths",
            "Realtime previews for faster stakeholder review",
            "A more coherent handoff between creation and deployment",
        ],
    },
    {
        "kind": "stack",
        "eyebrow": "04  Stack",
        "title": "A six-layer product stack",
        "layers": [
            ("01", "Editor cloud", "Scene composition, hierarchy control, materials, and environment setup in the browser."),
            ("02", "Realtime engine", "WebGPU and WebGL delivery, rendering control, interaction logic, and optimized runtime performance."),
            ("03", "Asset graph", "Models, textures, media, scripts, and structured reuse across product teams."),
            ("04", "Team operations", "Approvals, comments, review checkpoints, and shared production governance."),
            ("05", "Launch fabric", "Deployment to landing pages, product showcases, immersive demos, and XR endpoints."),
            ("06", "Data connectors", "Commerce, CMS, analytics, telemetry, and external system integration."),
        ],
    },
    {
        "kind": "mosaic",
        "eyebrow": "05  Surfaces",
        "title": "The platform is designed for multiple commercial surfaces, not one narrow use case.",
        "subtitle": "Loopnode can support browser games, XR experiences, product configurators, interactive campaigns, and other premium spatial interfaces.",
        "images": [IMG_GAMES, IMG_AR, IMG_VR, IMG_CONFIG, IMG_ADS],
    },
    {
        "kind": "signal_band",
        "eyebrow": "06  Feature pillars",
        "title": "The product combines creative speed with a more modern technical base.",
        "items": [
            ("WebGPU ready", "A modern graphics path while maintaining practical browser coverage."),
            ("Zero compile loops", "Faster iteration with a JavaScript-native workflow."),
            ("AI-assisted creation", "Vibe coding, editor automation, and natural-language acceleration."),
            ("Collaboration", "Realtime teamwork without turning review into a separate product."),
        ],
    },
    {
        "kind": "image_split_reverse",
        "eyebrow": "07  Buyers",
        "title": "Initial demand comes from teams building visually ambitious digital products.",
        "body": "The strongest wedge is not only gaming. Loopnode can be sold into studios, product teams, immersive commerce groups, simulation builders, and agencies shipping high-fidelity interactive work.",
        "image": IMG_ADS,
        "bullets": [
            "Game and interactive studios",
            "Commerce and configuration teams",
            "XR and simulation product groups",
            "Agencies building launch-ready immersive surfaces",
        ],
    },
    {
        "kind": "compare",
        "eyebrow": "08  Positioning",
        "title": "Where Loopnode is different",
        "left_title": "Loopnode",
        "right_title": "Typical stack",
        "rows": [
            ("Scene workflow", "Browser-native scene workflows", "Usually external or fragmented"),
            ("Review layer", "Review inside product flow", "Often handled with patchwork tools"),
            ("Product framing", "Launch-ready product framing", "Often engine-centric rather than product-centric"),
            ("Operational stack", "Connected runtime + asset + team layer", "Commonly assembled by hand"),
        ],
    },
    {
        "kind": "roadmap",
        "eyebrow": "09  Rollout",
        "title": "A practical expansion path",
        "steps": [
            ("Phase 1", "Launch editor, runtime, and browser publishing around a premium 3D creation wedge."),
            ("Phase 2", "Deepen asset operations, collaboration, and review workflows for production teams."),
            ("Phase 3", "Expand into AI-assisted creation, automation, and more structured deployment tooling."),
            ("Phase 4", "Broaden into more complex spatial product systems, data integrations, and enterprise surfaces."),
        ],
    },
    {
        "kind": "economics",
        "eyebrow": "10  Revenue model",
        "title": "Revenue can compound across software, workflow depth, and delivery infrastructure.",
        "columns": [
            ("Platform access", "Core subscription or enterprise access to the editing and runtime environment."),
            ("Production modules", "Higher-value layers for collaboration, governance, publishing, and analytics."),
            ("Services + integrations", "Implementation, migration, deployment configuration, and data connections."),
        ],
    },
    {
        "kind": "dark_thesis",
        "eyebrow": "11  Why this wins",
        "title": "The opportunity is to make browser 3D feel less like raw tooling and more like a finished product system.",
        "body": "Loopnode is differentiated by product coherence. Instead of asking teams to assemble scene tooling, runtime logic, review workflows, and launch mechanics on their own, it offers a cleaner operating layer designed around the actual shipping process.",
    },
    {
        "kind": "closing",
        "eyebrow": "12  Closing",
        "title": "Loopnode",
        "subtitle": "A premium browser-native 3D platform for teams building interactive products, immersive launches, and realtime spatial experiences.",
        "contact": "helpdesk@loopnode.net",
        "image": IMG_VR,
    },
]


def add_rect(slide, left, top, width, height, fill, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False, font="Arial", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_fit_text(slide, left, top, width, height, text, size, color, bold=False, font="Arial", align=PP_ALIGN.LEFT):
    box = add_text(slide, left, top, width, height, text, size=size, color=color, bold=bold, font=font, align=align)
    box.text_frame.fit_text(font_family=font, max_size=size, bold=bold)
    return box


def add_label(slide, text, dark=False):
    color = CYAN if dark else BLUE
    add_text(slide, Inches(0.7), Inches(0.45), Inches(4.0), Inches(0.35), text.upper(), size=12, color=color, bold=True, font="Arial")


def add_page_num(slide, num, dark=False):
    color = rgb("#6B7A8E") if not dark else rgb("#8EA7C4")
    add_text(slide, Inches(12.2), Inches(0.4), Inches(0.4), Inches(0.3), f"{num:02d}", size=11, color=color, font="Arial", align=PP_ALIGN.RIGHT)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def picture(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def cover(slide, data, idx):
    set_bg(slide, PAPER)
    add_page_num(slide, idx)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY, NAVY, radius=False)
    add_rect(slide, Inches(0.55), Inches(0.55), Inches(5.2), Inches(6.4), WHITE)
    add_label(slide, data["eyebrow"], dark=True)
    picture(slide, data["image"], Inches(6.0), Inches(0.0), Inches(7.333), Inches(7.5))
    if LOGO.exists():
        picture(slide, LOGO, Inches(0.95), Inches(0.95), Inches(1.2), Inches(1.2))
    add_fit_text(slide, Inches(0.95), Inches(2.0), Inches(4.2), Inches(1.9), data["title"], 25, INK, bold=True)
    add_text(slide, Inches(0.95), Inches(4.2), Inches(4.1), Inches(1.3), data["subtitle"], size=13.5, color=MUTED)
    add_rect(slide, Inches(0.95), Inches(5.75), Inches(3.35), Inches(0.62), BLUE, BLUE)
    add_text(slide, Inches(1.18), Inches(5.93), Inches(2.8), Inches(0.22), "Premium web-native 3D stack", size=12, color=WHITE, bold=True)


def statement(slide, data, idx):
    set_bg(slide, WHITE)
    add_page_num(slide, idx)
    add_label(slide, data["eyebrow"])
    add_fit_text(slide, Inches(0.82), Inches(1.0), Inches(8.0), Inches(1.4), data["title"], 28, INK, bold=True)
    add_text(slide, Inches(0.82), Inches(2.65), Inches(5.8), Inches(1.5), data["body"], size=16, color=MUTED)
    add_rect(slide, Inches(8.85), Inches(0.8), Inches(3.7), Inches(5.9), PAPER, LINE)
    y = 1.4
    for i, point in enumerate(data["points"], start=1):
        add_rect(slide, Inches(9.2), Inches(y), Inches(0.6), Inches(0.48), BLUE, BLUE)
        add_text(slide, Inches(9.42), Inches(y + 0.12), Inches(0.18), Inches(0.18), str(i), size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(10.0), Inches(y - 0.02), Inches(2.1), Inches(0.55), point, size=16, color=INK, bold=True)
        y += 1.45


def three_cards(slide, data, idx):
    set_bg(slide, PAPER)
    add_page_num(slide, idx)
    add_label(slide, data["eyebrow"])
    add_fit_text(slide, Inches(0.82), Inches(0.95), Inches(6.2), Inches(0.9), data["title"], 24, INK, bold=True)
    add_text(slide, Inches(0.82), Inches(1.85), Inches(7.0), Inches(0.6), data["subtitle"], size=14, color=MUTED)
    x = 0.82
    fills = [WHITE, WHITE, WHITE]
    accents = [CYAN, AMBER, VIOLET]
    for i, (title, body) in enumerate(data["cards"]):
        add_rect(slide, Inches(x), Inches(3.0), Inches(3.9), Inches(3.4), fills[i], LINE)
        add_rect(slide, Inches(x + 0.25), Inches(3.28), Inches(0.62), Inches(0.1), accents[i], accents[i], radius=False)
        add_text(slide, Inches(x + 0.25), Inches(3.55), Inches(3.2), Inches(0.6), title, size=17, color=INK, bold=True)
        add_text(slide, Inches(x + 0.25), Inches(4.35), Inches(3.15), Inches(1.5), body, size=13.5, color=MUTED)
        x += 4.15


def image_split(slide, data, idx, reverse=False):
    set_bg(slide, WHITE)
    add_page_num(slide, idx)
    add_label(slide, data["eyebrow"])
    if reverse:
        picture(slide, data["image"], Inches(0.7), Inches(1.05), Inches(5.1), Inches(5.75))
        tx = 6.2
    else:
        picture(slide, data["image"], Inches(7.15), Inches(1.05), Inches(5.45), Inches(5.75))
        tx = 0.78
    add_fit_text(slide, Inches(tx), Inches(1.15), Inches(5.3), Inches(1.0), data["title"], 24, INK, bold=True)
    add_text(slide, Inches(tx), Inches(2.25), Inches(4.95), Inches(1.05), data["body"], size=14.5, color=MUTED)
    y = 3.55
    for bullet in data["bullets"]:
        add_rect(slide, Inches(tx), Inches(y + 0.08), Inches(0.18), Inches(0.18), BLUE, BLUE, radius=False)
        add_text(slide, Inches(tx + 0.35), Inches(y), Inches(4.55), Inches(0.45), bullet, size=13.5, color=INK)
        y += 0.75


def stack(slide, data, idx):
    set_bg(slide, NAVY)
    add_page_num(slide, idx, dark=True)
    add_label(slide, data["eyebrow"], dark=True)
    add_fit_text(slide, Inches(0.78), Inches(1.0), Inches(5.0), Inches(0.8), data["title"], 24, WHITE, bold=True)
    y = 1.95
    fills = [SLATE, rgb("#162941"), rgb("#1A314D"), rgb("#1F3657"), rgb("#27426A"), rgb("#2E4C78")]
    for i, (num, title, body) in enumerate(data["layers"]):
        left = 0.78 if i < 3 else 6.75
        top = y + (i % 3) * 1.62
        add_rect(slide, Inches(left), Inches(top), Inches(5.75), Inches(1.28), fills[i], DARK_LINE)
        add_text(slide, Inches(left + 0.25), Inches(top + 0.2), Inches(0.5), Inches(0.35), num, size=12, color=LIME, bold=True)
        add_text(slide, Inches(left + 0.8), Inches(top + 0.15), Inches(2.8), Inches(0.35), title, size=16, color=WHITE, bold=True)
        add_text(slide, Inches(left + 0.8), Inches(top + 0.55), Inches(4.55), Inches(0.55), body, size=11.8, color=rgb("#C6D5E4"))


def mosaic(slide, data, idx):
    set_bg(slide, PAPER)
    add_page_num(slide, idx)
    add_label(slide, data["eyebrow"])
    add_fit_text(slide, Inches(0.82), Inches(0.92), Inches(6.6), Inches(0.9), data["title"], 24, INK, bold=True)
    add_text(slide, Inches(0.82), Inches(1.82), Inches(6.7), Inches(0.6), data["subtitle"], size=14, color=MUTED)
    picture(slide, data["images"][0], Inches(0.82), Inches(2.75), Inches(4.1), Inches(3.95))
    picture(slide, data["images"][1], Inches(5.1), Inches(2.75), Inches(2.25), Inches(1.85))
    picture(slide, data["images"][2], Inches(7.55), Inches(2.75), Inches(2.25), Inches(1.85))
    picture(slide, data["images"][3], Inches(5.1), Inches(4.85), Inches(2.25), Inches(1.85))
    picture(slide, data["images"][4], Inches(7.55), Inches(4.85), Inches(4.25), Inches(1.85))
    add_rect(slide, Inches(10.05), Inches(2.75), Inches(2.6), Inches(3.95), WHITE, LINE)
    add_text(slide, Inches(10.35), Inches(3.15), Inches(1.9), Inches(0.35), "Surface range", size=16, color=INK, bold=True)
    add_text(slide, Inches(10.35), Inches(3.7), Inches(1.95), Inches(2.0), "Games\nAR\nVR\nConfigurations\nInteractive ads", size=17, color=BLUE, bold=True)


def signal_band(slide, data, idx):
    set_bg(slide, WHITE)
    add_page_num(slide, idx)
    add_label(slide, data["eyebrow"])
    add_fit_text(slide, Inches(0.82), Inches(0.92), Inches(8.2), Inches(0.9), data["title"], 24, INK, bold=True)
    colors = [CYAN, AMBER, LIME, VIOLET]
    x = 0.82
    for i, (title, body) in enumerate(data["items"]):
        add_rect(slide, Inches(x), Inches(2.2), Inches(2.95), Inches(3.55), NAVY, DARK_LINE)
        add_rect(slide, Inches(x + 0.28), Inches(2.48), Inches(0.72), Inches(0.72), colors[i], colors[i])
        add_text(slide, Inches(x + 0.28), Inches(3.45), Inches(2.3), Inches(0.52), title, size=16, color=WHITE, bold=True)
        add_text(slide, Inches(x + 0.28), Inches(4.15), Inches(2.35), Inches(1.2), body, size=12.5, color=rgb("#C6D5E4"))
        x += 3.12


def compare(slide, data, idx):
    set_bg(slide, PAPER)
    add_page_num(slide, idx)
    add_label(slide, data["eyebrow"])
    add_fit_text(slide, Inches(0.82), Inches(0.95), Inches(6.6), Inches(0.8), data["title"], 24, INK, bold=True)
    add_text(slide, Inches(0.82), Inches(1.78), Inches(6.5), Inches(0.5), "A product-first comparison of how Loopnode is positioned against the usual assembled stack.", size=13.5, color=MUTED)

    add_rect(slide, Inches(0.82), Inches(2.35), Inches(5.58), Inches(0.72), NAVY, NAVY)
    add_rect(slide, Inches(6.93), Inches(2.35), Inches(5.58), Inches(0.72), rgb("#EDF2F8"), LINE)
    add_text(slide, Inches(1.05), Inches(2.57), Inches(4.8), Inches(0.25), data["left_title"], size=13, color=WHITE, bold=True)
    add_text(slide, Inches(7.18), Inches(2.57), Inches(4.8), Inches(0.25), data["right_title"], size=13, color=INK, bold=True)

    y = 3.28
    row_colors = [CYAN, AMBER, VIOLET, LIME]
    for label, left, right in data["rows"]:
        accent = row_colors[(int((y - 3.28) / 0.92)) % len(row_colors)]
        add_rect(slide, Inches(0.82), Inches(y), Inches(12.1), Inches(0.86), WHITE, LINE)
        add_rect(slide, Inches(1.02), Inches(y + 0.18), Inches(0.9), Inches(0.48), accent, accent)
        add_text(slide, Inches(1.18), Inches(y + 0.32), Inches(0.58), Inches(0.16), label, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(2.2), Inches(y + 0.18), Inches(4.0), Inches(0.42), left, size=13.8, color=BLUE, bold=True)
        add_text(slide, Inches(7.08), Inches(y + 0.18), Inches(4.75), Inches(0.42), right, size=13.2, color=MUTED)
        y += 0.98


def roadmap(slide, data, idx):
    set_bg(slide, WHITE)
    add_page_num(slide, idx)
    add_label(slide, data["eyebrow"])
    add_fit_text(slide, Inches(0.82), Inches(0.95), Inches(5.7), Inches(0.8), data["title"], 24, INK, bold=True)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.2), Inches(3.55), Inches(10.8), Inches(0.03)).fill.solid()
    line = slide.shapes[-1]
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    xs = [1.3, 4.05, 6.8, 9.55]
    colors = [BLUE, CYAN, AMBER, VIOLET]
    for i, (phase, text) in enumerate(data["steps"]):
        add_rect(slide, Inches(xs[i]), Inches(2.55), Inches(2.2), Inches(2.2), PAPER, LINE)
        add_rect(slide, Inches(xs[i] + 0.72), Inches(3.2), Inches(0.76), Inches(0.76), colors[i], colors[i])
        add_text(slide, Inches(xs[i] + 0.35), Inches(2.85), Inches(1.4), Inches(0.25), phase, size=12.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(xs[i] + 0.18), Inches(4.2), Inches(1.85), Inches(0.9), text, size=11.8, color=MUTED, align=PP_ALIGN.CENTER)


def economics(slide, data, idx):
    set_bg(slide, NAVY)
    add_page_num(slide, idx, dark=True)
    add_label(slide, data["eyebrow"], dark=True)
    add_fit_text(slide, Inches(0.82), Inches(0.95), Inches(8.0), Inches(0.8), data["title"], 24, WHITE, bold=True)
    x = 0.82
    colors = [CYAN, AMBER, LIME]
    for i, (title, body) in enumerate(data["columns"]):
        add_rect(slide, Inches(x), Inches(2.1), Inches(3.85), Inches(3.7), SLATE, DARK_LINE)
        add_rect(slide, Inches(x + 0.28), Inches(2.45), Inches(0.9), Inches(0.1), colors[i], colors[i], radius=False)
        add_text(slide, Inches(x + 0.28), Inches(2.8), Inches(2.8), Inches(0.4), title, size=17, color=WHITE, bold=True)
        add_text(slide, Inches(x + 0.28), Inches(3.55), Inches(3.0), Inches(1.6), body, size=13, color=rgb("#C6D5E4"))
        x += 4.08


def dark_thesis(slide, data, idx):
    set_bg(slide, INK)
    add_page_num(slide, idx, dark=True)
    add_label(slide, data["eyebrow"], dark=True)
    add_fit_text(slide, Inches(0.95), Inches(1.2), Inches(9.4), Inches(2.1), data["title"], 28, WHITE, bold=True)
    add_text(slide, Inches(0.95), Inches(4.15), Inches(7.9), Inches(1.6), data["body"], size=16, color=rgb("#C6D5E4"))
    add_rect(slide, Inches(9.7), Inches(1.5), Inches(2.2), Inches(3.9), WHITE, WHITE)
    if LOGO.exists():
        picture(slide, LOGO, Inches(10.15), Inches(2.15), Inches(1.3), Inches(1.3))
    add_text(slide, Inches(10.0), Inches(3.8), Inches(1.6), Inches(0.35), "Loopnode", size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)


def closing(slide, data, idx):
    set_bg(slide, PAPER)
    add_page_num(slide, idx)
    picture(slide, data["image"], Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY, NAVY, radius=False).fill.transparency = 0.28
    add_label(slide, data["eyebrow"], dark=True)
    if LOGO.exists():
        picture(slide, LOGO, Inches(0.95), Inches(1.2), Inches(1.2), Inches(1.2))
    add_fit_text(slide, Inches(0.95), Inches(2.05), Inches(5.2), Inches(0.9), data["title"], 30, WHITE, bold=True)
    add_text(slide, Inches(0.95), Inches(3.05), Inches(5.2), Inches(1.0), data["subtitle"], size=15, color=rgb("#D7E3EE"))
    add_rect(slide, Inches(0.95), Inches(5.2), Inches(3.0), Inches(0.72), WHITE, WHITE)
    add_text(slide, Inches(1.22), Inches(5.44), Inches(2.5), Inches(0.25), data["contact"], size=13, color=INK, bold=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(layout)
        kind = data["kind"]
        if kind == "cover":
            cover(slide, data, idx)
        elif kind == "statement":
            statement(slide, data, idx)
        elif kind == "three_cards":
            three_cards(slide, data, idx)
        elif kind == "image_split":
            image_split(slide, data, idx, reverse=False)
        elif kind == "image_split_reverse":
            image_split(slide, data, idx, reverse=True)
        elif kind == "stack":
            stack(slide, data, idx)
        elif kind == "mosaic":
            mosaic(slide, data, idx)
        elif kind == "signal_band":
            signal_band(slide, data, idx)
        elif kind == "compare":
            compare(slide, data, idx)
        elif kind == "roadmap":
            roadmap(slide, data, idx)
        elif kind == "economics":
            economics(slide, data, idx)
        elif kind == "dark_thesis":
            dark_thesis(slide, data, idx)
        elif kind == "closing":
            closing(slide, data, idx)
        else:
            raise ValueError(f"Unknown slide kind: {kind}")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
